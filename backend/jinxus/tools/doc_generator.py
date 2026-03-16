"""DocGenerator - 문서 생성 도구

Word(.docx), PowerPoint(.pptx) 문서 생성:
- 제목, 본문, 표 삽입
- 템플릿 기반 보고서 생성
- 마크다운 → 워드 변환
"""
import json
import logging
from pathlib import Path
from typing import Any

from .base import JinxTool, ToolResult

logger = logging.getLogger(__name__)


class DocGenerator(JinxTool):
    """Word/PowerPoint 문서 생성 도구"""

    name = "doc_generator"
    description = "Word(.docx) 또는 PowerPoint(.pptx) 문서를 생성합니다. 보고서, 제안서, 프레젠테이션 등을 만들 수 있습니다."
    allowed_agents = ["JX_WRITER", "JX_ANALYST", "JX_OPS", "JS_PERSONA"]
    input_schema = {
        "type": "object",
        "properties": {
            "format": {
                "type": "string",
                "enum": ["docx", "pptx"],
                "description": "출력 형식: docx(워드), pptx(파워포인트)",
            },
            "output_path": {
                "type": "string",
                "description": "출력 파일 경로",
            },
            "title": {
                "type": "string",
                "description": "문서 제목",
            },
            "content": {
                "type": "string",
                "description": "문서 본문 (마크다운 형식 가능). 파워포인트의 경우 '---'로 슬라이드 구분",
            },
            "author": {
                "type": "string",
                "description": "작성자 이름 (기본: JINXUS)",
                "default": "JINXUS",
            },
        },
        "required": ["format", "output_path", "title", "content"],
    }

    async def run(self, input_data: Any) -> ToolResult:
        if isinstance(input_data, str):
            input_data = json.loads(input_data)

        fmt = input_data.get("format", "docx")
        output_path = input_data.get("output_path", "")
        title = input_data.get("title", "제목 없음")
        content = input_data.get("content", "")
        author = input_data.get("author", "JINXUS")

        if not output_path:
            return ToolResult(success=False, output=None, error="output_path가 필요합니다")

        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)

        try:
            if fmt == "docx":
                return await self._create_docx(out, title, content, author)
            elif fmt == "pptx":
                return await self._create_pptx(out, title, content, author)
            else:
                return ToolResult(success=False, output=None, error=f"미지원 형식: {fmt}")
        except ImportError as e:
            return ToolResult(success=False, output=None, error=f"패키지 미설치: {e}")
        except Exception as e:
            logger.error(f"[DocGenerator] 오류: {e}")
            return ToolResult(success=False, output=None, error=str(e))

    async def _create_docx(self, path: Path, title: str, content: str, author: str) -> ToolResult:
        from docx import Document
        from docx.shared import Pt

        doc = Document()
        doc.core_properties.author = author

        # 제목
        doc.add_heading(title, level=0)

        # 본문 파싱 (마크다운 간이 변환)
        for line in content.split("\n"):
            stripped = line.strip()
            if not stripped:
                doc.add_paragraph("")
            elif stripped.startswith("### "):
                doc.add_heading(stripped[4:], level=3)
            elif stripped.startswith("## "):
                doc.add_heading(stripped[3:], level=2)
            elif stripped.startswith("# "):
                doc.add_heading(stripped[2:], level=1)
            elif stripped.startswith("- ") or stripped.startswith("* "):
                doc.add_paragraph(stripped[2:], style="List Bullet")
            elif stripped.startswith("1. ") or stripped.startswith("2. ") or stripped.startswith("3. "):
                doc.add_paragraph(stripped[3:], style="List Number")
            elif stripped.startswith("> "):
                p = doc.add_paragraph(stripped[2:])
                p.style = doc.styles["Quote"] if "Quote" in [s.name for s in doc.styles] else p.style
            elif stripped.startswith("| "):
                # 간단한 표 처리 (첫 번째 | 행에서 컬럼 수 결정)
                cells = [c.strip() for c in stripped.split("|")[1:-1]]
                if not hasattr(self, "_current_table"):
                    self._current_table = doc.add_table(rows=0, cols=len(cells))
                    self._current_table.style = "Table Grid"
                if not all(c.replace("-", "") == "" for c in cells):  # 구분선 건너뛰기
                    row = self._current_table.add_row()
                    for i, cell in enumerate(cells):
                        if i < len(row.cells):
                            row.cells[i].text = cell
            else:
                if hasattr(self, "_current_table"):
                    delattr(self, "_current_table")
                doc.add_paragraph(stripped)

        if hasattr(self, "_current_table"):
            delattr(self, "_current_table")

        doc.save(str(path))
        return ToolResult(success=True, output=f"Word 문서 생성 완료: {path}")

    async def _create_pptx(self, path: Path, title: str, content: str, author: str) -> ToolResult:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.core_properties.author = author

        # 제목 슬라이드
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        if slide.placeholders[1]:
            slide.placeholders[1].text = f"작성: {author}"

        # '---'로 슬라이드 구분
        slides_content = content.split("---")
        for slide_text in slides_content:
            slide_text = slide_text.strip()
            if not slide_text:
                continue

            lines = slide_text.split("\n")
            slide_title = lines[0].lstrip("#").strip() if lines else "슬라이드"
            body = "\n".join(lines[1:]).strip()

            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_title
            if body and slide.placeholders[1]:
                slide.placeholders[1].text = body

        prs.save(str(path))
        return ToolResult(
            success=True,
            output=f"PowerPoint 문서 생성 완료: {path} ({len(prs.slides)}장 슬라이드)",
        )
